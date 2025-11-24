import os
import pickle
from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain.embeddings import HuggingFaceEmbeddings


# ---------------- DOCUMENT LOADERS ---------------- #

def load_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def load_pdf(file):
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text


def load_docx(file):
    doc = Document(file)
    return "\n".join([p.text for p in doc.paragraphs])


def load_files(uploaded_files):
    text = ""
    for f in uploaded_files:
        if f.name.endswith(".pptx"):
            text += load_pptx(f)
        elif f.name.endswith(".pdf"):
            text += load_pdf(f)
        elif f.name.endswith(".docx"):
            text += load_docx(f)
    return text


# ---------------- EMBEDDING + CACHING ---------------- #

CACHE_DIR = "cache"
os.makedirs(CACHE_DIR, exist_ok=True)


def cache_path(files_hash):
    return os.path.join(CACHE_DIR, f"{files_hash}.pkl")


def load_cache(files_hash):
    path = cache_path(files_hash)
    if os.path.exists(path):
        with open(path, "rb") as f:
            return pickle.load(f)
    return None


def save_cache(files_hash, vector_db):
    with open(cache_path(files_hash), "wb") as f:
        pickle.dump(vector_db, f)


def hash_files(uploaded_files):
    parts = [f"{f.name}-{f.size}" for f in uploaded_files]
    return str(abs(hash("|".join(parts))))


def create_vector_store(text):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=70,
    )
    chunks = splitter.split_text(text)

    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    vector_db = FAISS.from_texts(chunks, embedding=embeddings)

    return vector_db


# --------- SEMANTIC WEB EXTRA HELPERS -------- #

def make_rdf_triples(text_chunk):
    """Simple heuristic RDF triple generator."""
    lines = [l.strip() for l in text_chunk.split(".") if len(l.split()) > 2]
    triples = []

    for ln in lines:
        parts = ln.split(" ")
        if len(parts) >= 3:
            subject = parts[0]
            predicate = parts[1]
            obj = " ".join(parts[2:])
            triples.append((subject, predicate, obj))

    return triples


def owl_relation(subject, relation, obj):
    return f"{subject}  {relation}  {obj}"


def compute_similarity_scores(query, docs, embedding_model):
    q_emb = embedding_model.embed_query(query)
    scores = []

    for d in docs:
        d_emb = embedding_model.embed_query(d.page_content)
        sim = sum(q * dd for q, dd in zip(q_emb, d_emb))
        scores.append(sim)

    return scores

